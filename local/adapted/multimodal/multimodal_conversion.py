import sys
from gtts import gTTS
from transformers import pipeline, AutoTokenizer, AutoModelForSeq2SeqLM 
from PIL import Image, ImageDraw, ImageFont
from moviepy.editor import VideoFileClip, AudioFileClip
from pptx import Presentation
import cv2
import numpy as np
import logging
import os
import tempfile
import textwrap
import requests
from io import BytesIO
from openai import OpenAI

# Set up logging
logging.basicConfig(level=logging.DEBUG)

# Set the model cache directory
os.environ['TRANSFORMERS_CACHE'] = '/app/multimodal/bart-large-cnn'

#Generate Audio
def convert_text_to_audio(text, audio_file):
    logging.debug("Starting audio conversion")
    tts = gTTS(text=text, lang='en')
    tts.save(audio_file)

#Generate Slides
def generate_slides_from_text(text, slides_file):
    logging.debug("Starting slides conversion")
    
    model_path = '/app/multimodal/bart-large-cnn'   
    
    try:
        logging.debug(f"Current working directory: {os.getcwd()}")
        logging.debug(f"Loading tokenizer from: {model_path}")
        tokenizer = AutoTokenizer.from_pretrained(model_path, local_files_only=True)
        logging.debug("Tokenizer loaded successfully")
        
        logging.debug(f"Loading model from: {model_path}")
        model = AutoModelForSeq2SeqLM.from_pretrained(model_path, local_files_only=True)
        logging.debug("Model loaded successfully")
        
        summarizer = pipeline("summarization", model=model, tokenizer=tokenizer)
        
        # Split the text into chunks
        max_chunk_length = 1000  
        text_chunks = textwrap.wrap(text, max_chunk_length)
        
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        
        for i, chunk in enumerate(text_chunks):
            logging.debug(f"Summarizing chunk {i+1}/{len(text_chunks)}")
            summary = summarizer(chunk, max_length=100, min_length=30, do_sample=False)[0]['summary_text']
            
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = f"Summary Slide {i+1}"
            content.text = summary
        
        prs.save(slides_file)
        logging.debug("Slides generated successfully")
        return [summarizer(chunk, max_length=100, min_length=30, do_sample=False)[0]['summary_text'] for chunk in text_chunks]
    except Exception as e:
        logging.error(f"Error in generate_slides: {str(e)}")
        raise
    
# Set up OpenAI API key for DALL-E
client = OpenAI(api_key =os.environ.get("OPENAI_API_KEY"))

#Generate Image
def generate_image(prompt):
    try:
        response = client.images.generate(
            model="dall-e-3",
            prompt=prompt,
            size="1024x1024",
            quality="standard",
            n=1,
        )
        logging.debug(f"Response structure: {response}")
        image_url = response.data[0].url
        response = requests.get(image_url)
        img = Image.open(BytesIO(response.content))
        return img
    except Exception as e:
        logging.error(f"Error generating image: {str(e)}")
        return None
    
#Generate Video    
def generate_video_from_text(text, audio_file, video_file):
    logging.debug("Starting video conversion")
    
    # Video settings
    width, height = 1280, 720
    fps = 24

    # Load audio file and get its duration
    audio = AudioFileClip(audio_file)
    duration = audio.duration

    # Create a temporary file
    with tempfile.NamedTemporaryFile(suffix=".mp4", delete=False) as temp_file:
        temp_video_file = temp_file.name

    # Create a VideoWriter object
    fourcc = cv2.VideoWriter_fourcc(*'mp4v')
    out = cv2.VideoWriter(temp_video_file, fourcc, fps, (width, height))

    # Create a black background image
    background = np.zeros((height, width, 3), dtype=np.uint8)

    try:
        font = ImageFont.truetype("arial.ttf", 50)
    except IOError:
        font = ImageFont.load_default()

    for frame in range(int(fps * duration)):
        # Wrap text to fit screen width
        wrapped_lines = textwrap.wrap(text, width=40)
        total_lines = len(wrapped_lines)
        line_duration = duration / total_lines

    # Generate an image for every 5 lines of text
    images = []
    for i in range(0, total_lines, 5):
        image_prompt = " ".join(wrapped_lines[i:i+5])
        img = generate_image(image_prompt)
        if img:
            img = img.resize((width // 2, height // 2))
            images.append(img)

    for frame in range(int(fps * duration)):
        img = Image.fromarray(background)
        draw = ImageDraw.Draw(img)

        current_time = frame / fps
        current_line_index = int(current_time / line_duration)

        # Display current line and next two lines
        y_position = height // 2 - 60
        for i in range(current_line_index, min(current_line_index + 3, total_lines)):
            line = wrapped_lines[i]
            draw.text((width // 4, y_position), line, font=font, fill=(255, 255, 255), anchor="mm")
            y_position += 60

        # Display the corresponding image
        image_index = current_line_index // 5
        if image_index < len(images):
            img.paste(images[image_index], (width // 2, height // 4))

        # Convert PIL Image to OpenCV format
        cv_img = cv2.cvtColor(np.array(img), cv2.COLOR_RGB2BGR)
        out.write(cv_img)

    out.release()

    try:
        # Combine video and audio using moviepy
        video = VideoFileClip(temp_video_file)
        video = video.set_audio(audio)
        video.write_videofile(video_file, codec="libx264", audio_codec="aac")
    finally:
        # Remove temporary video file
        os.remove(temp_video_file)

    logging.debug("Video generated successfully")
    
if __name__ == "__main__":
    logging.debug("Starting app.py")
    text = sys.argv[1]
    audio_file = sys.argv[2]
    slides_file = sys.argv[3]
    video_file = sys.argv[4]
    generate_audio = bool(int(sys.argv[5]))
    generate_slides_flag = bool(int(sys.argv[6])) 
    generate_video = bool(int(sys.argv[7]))

    if generate_audio:
        convert_text_to_audio(text, audio_file)
    
    if generate_slides_flag:
        generate_slides_from_text(text, slides_file)

    if generate_video:
        generate_video_from_text(text, audio_file, video_file)