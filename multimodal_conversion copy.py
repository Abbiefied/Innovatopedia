import sys
from gtts import gTTS
from transformers import pipeline, AutoTokenizer, AutoModelForSeq2SeqLM 
from pptx import Presentation
from pptx.util import Inches
import logging
import os

# Set up logging
logging.basicConfig(level=logging.DEBUG)

# Set the model cache directory
os.environ['TRANSFORMERS_CACHE'] = '/app/multimodal/bart-large-cnn'

def convert_text_to_audio(text, audio_file):
    logging.debug("Starting audio conversion")
    tts = gTTS(text=text, lang='en')
    tts.save(audio_file)


def generate_slides_from_text(text, slides_file):
    logging.debug("Starting slides conversion")
    
    model_path = '/app/multimodal/bart-large-cnn'
    
    try:
        logging.debug(f"Current working directory: {os.getcwd()}")
        logging.debug(f"Contents of /app: {os.listdir('/app')}")
        logging.debug(f"Contents of model path: {os.listdir(model_path)}")
        
        logging.debug(f"Loading tokenizer from: {model_path}")
        tokenizer = AutoTokenizer.from_pretrained(model_path, local_files_only=True)
        logging.debug("Tokenizer loaded successfully")
        
        logging.debug(f"Loading model from: {model_path}")
        model = AutoModelForSeq2SeqLM.from_pretrained(model_path, local_files_only=True)
        logging.debug("Model loaded successfully")
        
        summarizer = pipeline("summarization", model=model, tokenizer=tokenizer)
        
        summary = summarizer(text, max_length=100, min_length=30, do_sample=False)[0]['summary_text']
        
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = "Summary Slide"
        content.text = summary

        prs.save(slides_file)
        logging.debug("Slides generated successfully")
    except Exception as e:
        logging.error(f"Error in generate_slides: {str(e)}")
        logging.error(f"Current working directory: {os.getcwd()}")
        logging.error(f"Contents of /app: {os.listdir('/app')}")
        if os.path.exists('/app/model_cache'):
            logging.error(f"Contents of /app/model_cache: {os.listdir('/app/model_cache')}")
        raise

if __name__ == "__main__":
    logging.debug("Starting app.py")
    text = sys.argv[1]
    audio_file = sys.argv[2]
    slides_file = sys.argv[3]
    generate_audio = bool(int(sys.argv[4]))
    generate_slides_flag = bool(int(sys.argv[5])) 

    if generate_audio:
        convert_text_to_audio(text, audio_file)
    
    if generate_slides_flag:
        generate_slides_from_text(text, slides_file)
